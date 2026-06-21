"""Genera un PowerPoint de kickoff EPM con 7 diapositivas usando python-pptx.

Aplica el formato Deloitte definido en CLAUDE.md: fondo blanco, fuente Arial,
linea verde de header en cada slide, titulos en negro Deloitte y cuerpo en gris.

Pensado para ser invocado desde New-KickoffPPT.ps1, que valida los parametros
y construye la llamada. Tambien puede ejecutarse de forma aislada:

    python build_kickoff_ppt.py --client "ACME" --tool OneStream \
        --date 2026-06-21 --team-lead "Victor Argenta" --output "./Kickoff-ACME.pptx"

El contenido es una plantilla profesional con placeholders editables: no
inventa datos del cliente, solo deja la estructura lista para rellenar.
"""

import argparse
import sys

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# Paleta Deloitte (ver CLAUDE.md). Solo se usan los tonos necesarios aqui.
DELOITTE_BLACK = RGBColor(0x26, 0x28, 0x2A)   # titulares y texto fuerte
DELOITTE_DARK_GRAY = RGBColor(0x53, 0x56, 0x5A)  # subtitulos y contraste
DELOITTE_MID_GRAY = RGBColor(0x63, 0x66, 0x6A)   # body y captions
DELOITTE_GREEN = RGBColor(0x86, 0xBC, 0x25)   # acento unico
DELOITTE_DARK_GREEN = RGBColor(0x04, 0x6A, 0x38)  # headers, estados, links
DELOITTE_OFFWHITE = RGBColor(0xF6, 0xF6, 0xF6)   # fondo de tarjetas
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"

# Las 7 fases EPM, coherentes con New-EPMProject.ps1.
EPM_PHASES = [
    ("Kickoff", "Arranque del proyecto: alcance, equipo, planificacion."),
    ("Functional specs", "Especificaciones funcionales y requisitos."),
    ("Data model", "Diseno del modelo de datos y dimensiones."),
    ("Testing", "Casos de prueba, UAT y evidencias."),
    ("Training", "Material y sesiones de formacion."),
    ("Go-live", "Puesta en produccion y checklist de salida."),
    ("Hypercare", "Soporte intensivo posterior al go-live."),
]

# Etiqueta legible de la herramienta EPM para mostrar en las diapositivas.
TOOL_LABELS = {
    "OneStream": "OneStream",
    "OracleEPM": "Oracle EPM",
    "Anaplan": "Anaplan",
}


def _add_blank_slide(prs):
    """Devuelve una diapositiva en blanco con fondo blanco explicito."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def _add_green_header_line(slide, slide_width):
    """Linea verde Deloitte (3pt) en la parte superior de la slide."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), slide_width - Inches(1.2), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = DELOITTE_GREEN
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def _style_run(run, size, color, bold=False):
    """Aplica fuente Arial, tamano, color y bold a un run."""
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_content_title(slide, text):
    """Titulo de contenido: bold 20pt, negro Deloitte."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    _style_run(run, 20, DELOITTE_BLACK, bold=True)
    return box


def _add_bullets(slide, items, top=1.7, left=0.8, width=11.7, height=5.3):
    """Anade vinetas a partir de tuplas (nivel, texto).

    Nivel 0 = subtitulo de seccion (negro Deloitte, bold); nivel 1 = body gris.
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (level, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = text
        if level == 0:
            _style_run(run, 18, DELOITTE_BLACK, bold=True)
        else:
            _style_run(run, 15, DELOITTE_MID_GRAY)
    return box


def build_cover(prs, client, tool_label, date, team_lead):
    """Diapositiva 1: portada. Fondo blanco, linea verde, titulo 28pt negro."""
    slide = _add_blank_slide(prs)
    _add_green_header_line(slide, prs.slide_width)

    # Titulo principal de portada: bold 28pt negro Deloitte.
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.2))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    run = ttf.paragraphs[0].add_run()
    run.text = "Kickoff de Proyecto EPM"
    _style_run(run, 28, DELOITTE_BLACK, bold=True)

    # Nombre de cliente como subtitulo en gris oscuro de contraste.
    sub = ttf.add_paragraph()
    sub_run = sub.add_run()
    sub_run.text = client
    _style_run(sub_run, 20, DELOITTE_DARK_GRAY, bold=True)

    # Metadatos en gris medio (body/captions).
    meta_lines = [f"Herramienta EPM: {tool_label}", f"Fecha: {date}"]
    if team_lead:
        meta_lines.append(f"Team Lead: {team_lead}")
    box = slide.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.6))
    mtf = box.text_frame
    mtf.word_wrap = True
    for i, line in enumerate(meta_lines):
        p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        _style_run(run, 14, DELOITTE_MID_GRAY)


def build_agenda(prs):
    """Diapositiva 2: agenda (refleja las propias secciones de la sesion)."""
    slide = _add_blank_slide(prs)
    _add_green_header_line(slide, prs.slide_width)
    _add_content_title(slide, "Agenda")
    items = [
        (0, "1. Objetivos del proyecto"),
        (0, "2. Alcance"),
        (0, "3. Equipo y roles"),
        (0, "4. Hitos y timeline"),
        (0, "5. Proximos pasos"),
    ]
    _add_bullets(slide, items)


def build_objectives(prs):
    """Diapositiva 3: objetivos (placeholders editables)."""
    slide = _add_blank_slide(prs)
    _add_green_header_line(slide, prs.slide_width)
    _add_content_title(slide, "Objetivos")
    items = [
        (0, "Objetivo de negocio"),
        (1, "[Describir el resultado de negocio esperado]"),
        (0, "Objetivos del proyecto"),
        (1, "[Objetivo 1]"),
        (1, "[Objetivo 2]"),
        (1, "[Objetivo 3]"),
        (0, "Criterios de exito"),
        (1, "[Como mediremos que el proyecto ha tenido exito]"),
    ]
    _add_bullets(slide, items)


def build_scope(prs):
    """Diapositiva 4: alcance (in-scope / out-of-scope)."""
    slide = _add_blank_slide(prs)
    _add_green_header_line(slide, prs.slide_width)
    _add_content_title(slide, "Alcance")
    items = [
        (0, "Dentro de alcance"),
        (1, "[Procesos / modulos incluidos]"),
        (1, "[Entidades / dimensiones]"),
        (1, "[Integraciones]"),
        (0, "Fuera de alcance"),
        (1, "[Lo que explicitamente NO se aborda]"),
        (0, "Supuestos y dependencias"),
        (1, "[Supuestos clave del proyecto]"),
    ]
    _add_bullets(slide, items)


def build_team(prs, team_lead):
    """Diapositiva 5: equipo y roles."""
    slide = _add_blank_slide(prs)
    _add_green_header_line(slide, prs.slide_width)
    _add_content_title(slide, "Equipo y roles")
    lead_text = team_lead if team_lead else "[Nombre]"
    items = [
        (0, f"Team Lead: {lead_text}"),
        (1, "Responsable de la entrega y punto de contacto principal."),
        (0, "Sponsor de cliente"),
        (1, "[Nombre / rol]"),
        (0, "Consultores funcionales"),
        (1, "[Nombres]"),
        (0, "Consultores tecnicos"),
        (1, "[Nombres]"),
        (0, "Key users de cliente"),
        (1, "[Nombres por area]"),
    ]
    _add_bullets(slide, items)


def build_timeline(prs):
    """Diapositiva 6: hitos / timeline a partir de las 7 fases EPM."""
    slide = _add_blank_slide(prs)
    _add_green_header_line(slide, prs.slide_width)
    _add_content_title(slide, "Hitos y timeline")
    items = []
    for idx, (name, purpose) in enumerate(EPM_PHASES, start=1):
        items.append((0, f"Fase {idx}: {name}"))
        items.append((1, f"{purpose} [Fechas: por definir]"))
    _add_bullets(slide, items, height=5.6)


def build_next_steps(prs):
    """Diapositiva 7: proximos pasos."""
    slide = _add_blank_slide(prs)
    _add_green_header_line(slide, prs.slide_width)
    _add_content_title(slide, "Proximos pasos")
    items = [
        (0, "Acciones inmediatas"),
        (1, "[Accion] - Owner: [nombre] - Fecha: [dd/mm]"),
        (1, "[Accion] - Owner: [nombre] - Fecha: [dd/mm]"),
        (0, "Decisiones pendientes"),
        (1, "[Decision que requiere cierre]"),
        (0, "Proxima reunion"),
        (1, "[Fecha y objetivo del siguiente checkpoint]"),
    ]
    _add_bullets(slide, items)


def build_presentation(client, tool, date, team_lead, output):
    """Construye y guarda la presentacion completa."""
    tool_label = TOOL_LABELS.get(tool, tool)

    prs = Presentation()
    # Formato panoramico 16:9.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover(prs, client, tool_label, date, team_lead)
    build_agenda(prs)
    build_objectives(prs)
    build_scope(prs)
    build_team(prs, team_lead)
    build_timeline(prs)
    build_next_steps(prs)

    prs.save(output)
    return output


def main(argv=None):
    parser = argparse.ArgumentParser(
        description="Genera un PowerPoint de kickoff EPM (7 diapositivas)."
    )
    parser.add_argument("--client", required=True, help="Nombre del cliente.")
    parser.add_argument(
        "--tool",
        default="OneStream",
        choices=["OneStream", "OracleEPM", "Anaplan"],
        help="Herramienta EPM.",
    )
    parser.add_argument("--date", required=True, help="Fecha del proyecto (yyyy-MM-dd).")
    parser.add_argument("--team-lead", default="", help="Team Lead del proyecto.")
    parser.add_argument("--output", required=True, help="Ruta del .pptx de salida.")
    args = parser.parse_args(argv)

    try:
        path = build_presentation(
            args.client, args.tool, args.date, args.team_lead, args.output
        )
    except Exception as exc:  # noqa: BLE001 - reportar cualquier fallo al wrapper
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1

    print(path)
    return 0


if __name__ == "__main__":
    sys.exit(main())
